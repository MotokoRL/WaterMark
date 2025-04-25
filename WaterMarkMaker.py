import tkinter as tk
from tkinter import filedialog, messagebox
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.shared import RGBColor
from pptx import Presentation
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.colors import Color
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import io
import os

class WatermarkApp:
    def __init__(self, root):
        self.root = root
        self.root.title("水印添加程序")
        
        # 文件选择
        tk.Label(root, text="输入文件:").grid(row=0, column=0)
        self.input_file_entry = tk.Entry(root, width=40)
        self.input_file_entry.grid(row=0, column=1)
        tk.Button(root, text="选择文件", command=self.select_input_file).grid(row=0, column=2)
        
        # 保存路径
        tk.Label(root, text="保存路径:").grid(row=1, column=0)
        self.output_file_entry = tk.Entry(root, width=40)
        self.output_file_entry.grid(row=1, column=1)
        tk.Button(root, text="选择路径", command=self.select_output_file).grid(row=1, column=2)
        
        # 水印内容
        tk.Label(root, text="水印内容:").grid(row=2, column=0)
        self.watermark_text_entry = tk.Entry(root, width=40)
        self.watermark_text_entry.grid(row=2, column=1)
        
        # 字体选择
        tk.Label(root, text="字体文件:").grid(row=3, column=0)
        self.font_path_entry = tk.Entry(root, width=40)
        self.font_path_entry.grid(row=3, column=1)
        tk.Button(root, text="选择字体", command=self.select_font_file).grid(row=3, column=2)
        
        # 字体大小
        tk.Label(root, text="字体大小:").grid(row=4, column=0)
        self.font_size_entry = tk.Entry(root, width=40)
        self.font_size_entry.grid(row=4, column=1)
        
        # 字体颜色
        tk.Label(root, text="字体颜色 (RGB):").grid(row=5, column=0)
        self.color_entry = tk.Entry(root, width=40)
        self.color_entry.grid(row=5, column=1)
        self.color_entry.insert(0, "0,0,0")
        
        # 透明度
        tk.Label(root, text="透明度 (0.0-1.0):").grid(row=6, column=0)
        self.opacity_entry = tk.Entry(root, width=40)
        self.opacity_entry.grid(row=6, column=1)
        self.opacity_entry.insert(0, "0.5")
        
        # 密度
        tk.Label(root, text="密度 (1-10):").grid(row=7, column=0)
        self.density_entry = tk.Entry(root, width=40)
        self.density_entry.grid(row=7, column=1)
        self.density_entry.insert(0, "5")
        
        # 操作按钮
        tk.Button(root, text="添加水印", command=self.process_watermark).grid(row=8, column=1)
        tk.Button(root, text="退出", command=root.quit).grid(row=8, column=2)

    def select_input_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("所有文件", "*.*")])
        if file_path:
            self.input_file_entry.delete(0, tk.END)
            self.input_file_entry.insert(0, file_path)

    def select_output_file(self):
        file_path = filedialog.asksaveasfilename(defaultextension=".*")
        if file_path:
            self.output_file_entry.delete(0, tk.END)
            self.output_file_entry.insert(0, file_path)

    def select_font_file(self):
        font_path = filedialog.askopenfilename(filetypes=[("TrueType 字体", "*.ttf")])
        if font_path:
            self.font_path_entry.delete(0, tk.END)
            self.font_path_entry.insert(0, font_path)

    def validate_inputs(self):
        try:
            # 验证输入文件
            input_file = self.input_file_entry.get()
            if not input_file or not os.path.exists(input_file):
                messagebox.showerror("错误", "请选择有效的输入文件！")
                return False

            # 验证保存路径
            output_file = self.output_file_entry.get()
            if not output_file:
                messagebox.showerror("错误", "请选择保存路径！")
                return False

            # 验证水印内容
            watermark_text = self.watermark_text_entry.get()
            if not watermark_text:
                messagebox.showerror("错误", "水印内容不能为空！")
                return False

            # 验证字体文件
            font_path = self.font_path_entry.get()
            if not font_path or not os.path.exists(font_path):
                messagebox.showerror("错误", "请选择有效的字体文件！")
                return False

            # 验证字体大小
            font_size = int(self.font_size_entry.get())
            if font_size <= 0:
                messagebox.showerror("错误", "字体大小必须大于0！")
                return False

            # 验证颜色
            color = tuple(map(int, self.color_entry.get().split(',')))
            if len(color) != 3 or not all(0 <= c <= 255 for c in color):
                messagebox.showerror("错误", "颜色格式不正确！")
                return False

            # 验证透明度
            opacity = float(self.opacity_entry.get())
            if not 0.0 <= opacity <= 1.0:
                messagebox.showerror("错误", "透明度必须在0.0到1.0之间！")
                return False

            # 验证密度
            density = int(self.density_entry.get())
            if not 1 <= density <= 10:
                messagebox.showerror("错误", "密度必须在1到10之间！")
                return False

            return True
        except ValueError as e:
            messagebox.showerror("输入错误", f"输入格式不正确：{str(e)}")
            return False

    def process_watermark(self):
        if not self.validate_inputs():
            return

        watermark_details = {
            "watermark_text": self.watermark_text_entry.get(),
            "font_path": self.font_path_entry.get(),
            "font_size": int(self.font_size_entry.get()),
            "color": tuple(map(int, self.color_entry.get().split(','))),
            "opacity": float(self.opacity_entry.get()),
            "density": int(self.density_entry.get())
        }

        input_file = self.input_file_entry.get()
        output_file = self.output_file_entry.get()

        try:
            if input_file.endswith(".docx"):
                self.add_watermark_to_word(input_file, output_file, watermark_details)
            elif input_file.endswith(".pptx"):
                self.add_watermark_to_pptx(input_file, output_file, watermark_details)
            elif input_file.endswith(".pdf"):
                self.add_watermark_to_pdf(input_file, output_file, watermark_details)
            else:
                messagebox.showwarning("格式错误", "不支持的文件格式！")
        except Exception as e:
            messagebox.showerror("错误", f"处理文件时出错：{str(e)}")

    def add_watermark_to_word(self, input_path, output_path, watermark_details):
        doc = Document(input_path)
        for section in doc.sections:
            header = section.header
            paragraph = header.paragraphs[0]
            run = paragraph.add_run(watermark_details["watermark_text"])
            run.font.size = watermark_details["font_size"]
            run.font.name = watermark_details["font_path"]
            run.font.color.rgb = RGBColor(*watermark_details["color"][:3])
        doc.save(output_path)
        messagebox.showinfo("成功", "Word文件水印添加成功！")

    def add_watermark_to_pptx(self, input_path, output_path, watermark_details):
        prs = Presentation(input_path)
        for slide in prs.slides:
            left = top = 0
            width = height = prs.slide_width
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = watermark_details["watermark_text"]
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = watermark_details["font_size"]
                    run.font.name = watermark_details["font_path"]
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = Color(*watermark_details["color"], alpha=watermark_details["opacity"])
        prs.save(output_path)
        messagebox.showinfo("成功", "PowerPoint文件水印添加成功！")

    def add_watermark_to_pdf(self, input_path, output_path, watermark_details):
        output_pdf = PdfWriter()
        with open(input_path, "rb") as input_pdf:
            reader = PdfReader(input_pdf)
            total_pages = len(reader.pages)
            messagebox.showinfo("进度", f"开始处理PDF文件，共 {total_pages} 页...")
            for page in reader.pages:
                packet = io.BytesIO()
                c = canvas.Canvas(packet, pagesize=letter)
                font_name = "CustomFont_" + watermark_details["font_path"].replace("/", "_")
                try:
                    pdfmetrics.registerFont(TTFont(font_name, watermark_details["font_path"]))
                except:
                    font_name = "Helvetica"
                c.setFont(font_name, watermark_details["font_size"]
                c.setFillColor(Color(*watermark_details["color"], alpha=watermark_details["opacity"]))
                
                # Apply density by repeating the watermark
                for i in range(watermark_details["density"]):
                    x = (i + 1) * (letter[0] / (watermark_details["density"] + 1))
                    y = (i + 1) * (letter[1] / (watermark_details["density"] + 1))
                    c.drawString(x, y, watermark_details["watermark_text"])
                
                c.save()
                packet.seek(0)
                new_pdf = PdfReader(packet)
                page.merge_page(new_pdf.pages[0])
                output_pdf.add_page(page)
            
            with open(output_path, "wb") as output_file:
                output_pdf.write(output_file)
            messagebox.showinfo("成功", "PDF文件水印添加成功！")

def main():
    root = tk.Tk()
    app = WatermarkApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()